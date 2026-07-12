from __future__ import annotations

import re
import shutil
import subprocess
import uuid
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from .formulas import omml_text_to_latex
from .models import Document, Formula, ImageAsset, Page

OMML = "{http://schemas.openxmlformats.org/officeDocument/2006/math}"


def _safe_name(path: Path) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "_", path.stem).strip("_") or "document"


def _save_blob(blob: bytes, destination: Path, index: int, page: int, suffix: str) -> Path:
    destination.mkdir(parents=True, exist_ok=True)
    out = destination / f"page{page:03d}_img{index:02d}{suffix or '.bin'}"
    out.write_bytes(blob)
    return out


class DocumentExtractor:
    def extract(self, source: Path, asset_dir: Path) -> Document:
        ext = source.suffix.lower()
        if ext in {".ppt", ".doc"}:
            source = OfficeConverter().convert(source)
            ext = source.suffix.lower()
        if ext == ".pptx": return PptxExtractor().extract(source, asset_dir)
        if ext == ".docx": return DocxExtractor().extract(source, asset_dir)
        if ext == ".pdf": return PdfExtractor().extract(source, asset_dir)
        raise ValueError(f"不支持的文件格式: {ext}")


class PptxExtractor:
    def extract(self, source: Path, asset_dir: Path) -> Document:
        from pptx import Presentation
        presentation = Presentation(source)
        pages: list[Page] = []
        for page_no, slide in enumerate(presentation.slides, 1):
            page = Page(number=page_no)
            image_no = 0
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = shape.text.strip()
                    if value:
                        if not page.title:
                            page.title = value.splitlines()[0][:120]
                        page.text.append(value)
                if getattr(shape, "has_table", False):
                    page.tables.append([[cell.text.strip() for cell in row.cells] for row in shape.table.rows])
                if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        image_no += 1
                        suffix = Path(image.filename or "image.png").suffix
                        saved = _save_blob(image.blob, asset_dir, image_no, page_no, suffix)
                        page.images.append(ImageAsset(path=saved.name, source_page=page_no))
                    except ValueError:
                        # Linked images have no embeddable blob; their text is still retained.
                        pass
            try:
                notes = "\n".join(shape.text for shape in slide.notes_slide.shapes if getattr(shape, "has_text_frame", False)).strip()
                if notes:
                    page.text.append(f"讲者备注：\n{notes}")
            except Exception:
                pass
            pages.append(page)
        self._read_omml(source, pages)
        return Document(str(source), "pptx", pages)

    def _read_omml(self, source: Path, pages: list[Page]) -> None:
        with zipfile.ZipFile(source) as archive:
            for name in archive.namelist():
                match = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", name)
                if not match: continue
                root = ET.fromstring(archive.read(name))
                for math in root.findall(f".//{OMML}oMath") + root.findall(f".//{OMML}oMathPara"):
                    raw = "".join(math.itertext()).strip()
                    if raw:
                        pages[int(match.group(1)) - 1].formulas.append(Formula(omml_text_to_latex(raw), int(match.group(1)), "omml"))


class DocxExtractor:
    def extract(self, source: Path, asset_dir: Path) -> Document:
        from docx import Document as WordDocument
        document = WordDocument(source)
        page = Page(number=1)
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                if not page.title: page.title = text[:120]
                page.text.append(text)
        for table in document.tables:
            page.tables.append([[cell.text.strip() for cell in row.cells] for row in table.rows])
        with zipfile.ZipFile(source) as archive:
            image_no = 0
            for name in archive.namelist():
                if name.startswith("word/media/"):
                    image_no += 1
                    out = _save_blob(archive.read(name), asset_dir, image_no, 1, Path(name).suffix)
                    page.images.append(ImageAsset(out.name, 1))
                elif name == "word/document.xml":
                    root = ET.fromstring(archive.read(name))
                    for math in root.findall(f".//{OMML}oMath") + root.findall(f".//{OMML}oMathPara"):
                        raw = "".join(math.itertext()).strip()
                        if raw: page.formulas.append(Formula(omml_text_to_latex(raw), 1, "omml"))
        return Document(str(source), "docx", [page])


class PdfExtractor:
    def extract(self, source: Path, asset_dir: Path) -> Document:
        from pypdf import PdfReader
        reader = PdfReader(source)
        pages = []
        for number, pdf_page in enumerate(reader.pages, 1):
            text = (pdf_page.extract_text() or "").strip()
            page = Page(number=number, title=(text.splitlines()[0][:120] if text else f"第 {number} 页"))
            if text: page.text.append(text)
            image_no = 0
            try:
                for image in pdf_page.images:
                    image_no += 1
                    out = _save_blob(image.data, asset_dir, image_no, number, Path(image.name).suffix)
                    page.images.append(ImageAsset(out.name, number))
            except Exception:
                pass
            if len(text) < 40:
                try:
                    import pdfplumber
                    with pdfplumber.open(source) as pdf:
                        rendered = pdf.pages[number - 1].to_image(resolution=160)
                        out = asset_dir / f"page{number:03d}_render.png"
                        asset_dir.mkdir(parents=True, exist_ok=True)
                        rendered.save(out, format="PNG")
                        page.images.append(ImageAsset(out.name, number, "PDF 渲染页"))
                except Exception:
                    pass
            pages.append(page)
        return Document(str(source), "pdf", pages)


class OfficeConverter:
    """使用 VBS 调用已安装的 Office；仅在旧格式输入时使用。"""
    def convert(self, source: Path) -> Path:
        # Use the project-local work directory: managed Windows environments can
        # deny writes to the account-level temporary directory.
        work = Path.cwd() / "work"
        work.mkdir(parents=True, exist_ok=True)
        destination_dir = work / f"cmd_office_{uuid.uuid4().hex[:12]}"
        destination_dir.mkdir()
        staged = destination_dir / ("source" + source.suffix.lower())
        shutil.copy2(source, staged)
        output = destination_dir / ("converted.pptx" if source.suffix.lower() == ".ppt" else "converted.docx")
        app = "PowerPoint.Application" if source.suffix.lower() == ".ppt" else "Word.Application"
        # 24: ppSaveAsOpenXMLPresentation; 16: wdFormatDocumentDefault
        fmt = 24 if source.suffix.lower() == ".ppt" else 16
        if source.suffix.lower() == ".ppt":
            script = f'Set a=CreateObject("{app}"): Set d=a.Presentations.Open("{staged}",0,0,0): d.SaveAs "{output}",{fmt}: d.Close: a.Quit'
        else:
            script = f'Set a=CreateObject("{app}"): Set d=a.Documents.Open("{staged}"): d.SaveAs2 "{output}",{fmt}: d.Close: a.Quit'
        vbs = destination_dir / "convert.vbs"; vbs.write_text(script, encoding="utf-8")
        subprocess.run(["cscript.exe", "//nologo", str(vbs)], check=True, timeout=180)
        if not output.exists(): raise RuntimeError("Office 转换未生成输出文件")
        return output
